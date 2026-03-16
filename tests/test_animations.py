"""Tests for src/animations.py."""

import pytest
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from src.animations import (
    EFFECT_MAP,
    TARGET_SHAPE_MAP,
    _build_click_par,
    _build_visual_effect_xml,
    _get_shape_id,
    _resolve_target_shapes,
    apply_animations,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide(layout_index: int = 1) -> tuple:
    """Return (Presentation, Slide) for a standard layout."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    return prs, slide


# ---------------------------------------------------------------------------
# EFFECT_MAP & TARGET_SHAPE_MAP
# ---------------------------------------------------------------------------


def test_effect_map_has_known_effects():
    expected = {"appear", "fade", "fly-in", "wipe", "zoom", "float-in", "split", "blinds"}
    assert expected.issubset(set(EFFECT_MAP.keys()))


def test_effect_map_entries_have_required_keys():
    for name, entry in EFFECT_MAP.items():
        assert "preset" in entry, f"{name} missing 'preset'"
        assert "cls" in entry, f"{name} missing 'cls'"
        assert "subtype" in entry, f"{name} missing 'subtype'"
        assert "visual" in entry, f"{name} missing 'visual'"


def test_target_shape_map_has_standard_targets():
    for key in ("title", "content", "left", "right", "image", "subtitle"):
        assert key in TARGET_SHAPE_MAP


# ---------------------------------------------------------------------------
# _resolve_target_shapes
# ---------------------------------------------------------------------------


def test_resolve_target_shapes_title():
    _, slide = _make_slide(0)  # Title slide
    shapes = _resolve_target_shapes(slide, "title")
    assert len(shapes) >= 1
    assert any("title" in s.name.lower() for s in shapes)


def test_resolve_target_shapes_content():
    _, slide = _make_slide(1)  # Title + Content
    shapes = _resolve_target_shapes(slide, "content")
    assert len(shapes) >= 1


def test_resolve_target_shapes_all():
    _, slide = _make_slide(1)
    shapes = _resolve_target_shapes(slide, "all")
    assert len(shapes) == len(list(slide.shapes))


def test_resolve_target_shapes_bullets_alias():
    _, slide = _make_slide(1)
    shapes_content = _resolve_target_shapes(slide, "content")
    shapes_bullets = _resolve_target_shapes(slide, "bullets")
    assert shapes_content == shapes_bullets


def test_resolve_target_shapes_no_match():
    _, slide = _make_slide(1)
    shapes = _resolve_target_shapes(slide, "nonexistent_widget_xyz")
    assert shapes == []


# ---------------------------------------------------------------------------
# _get_shape_id
# ---------------------------------------------------------------------------


def test_get_shape_id_returns_string():
    _, slide = _make_slide(1)
    for shape in slide.shapes:
        sid = _get_shape_id(shape)
        assert isinstance(sid, str)
        assert sid.isdigit()


# ---------------------------------------------------------------------------
# _build_visual_effect_xml
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("visual", ["fade", "wipe", "zoom", "float", "split", "blinds"])
def test_build_visual_effect_xml_non_empty(visual):
    xml = _build_visual_effect_xml("5", visual)
    assert xml != ""
    assert 'spid="5"' in xml


def test_build_visual_effect_xml_none():
    xml = _build_visual_effect_xml("5", "none")
    assert xml == ""


@pytest.mark.parametrize("direction", ["bottom", "top", "left", "right"])
def test_build_visual_effect_xml_fly_directions(direction):
    xml = _build_visual_effect_xml("5", f"fly-{direction}")
    assert xml != ""
    assert "p:anim" in xml


# ---------------------------------------------------------------------------
# _build_click_par
# ---------------------------------------------------------------------------


def test_build_click_par_is_valid_xml():
    effect = EFFECT_MAP["fade"]
    xml = _build_click_par("3", effect)
    # Should parse as valid XML
    el = etree.fromstring(xml.encode())
    assert el.tag.endswith("}par") or el.tag == "p:par"


def test_build_click_par_contains_shape_id():
    effect = EFFECT_MAP["appear"]
    xml = _build_click_par("42", effect)
    assert 'spid="42"' in xml


# ---------------------------------------------------------------------------
# apply_animations – integration
# ---------------------------------------------------------------------------


def test_apply_animations_with_valid_effect():
    _, slide = _make_slide(1)
    animations = [{"target": "title", "effect": "fade"}]
    apply_animations(slide, animations)
    # Verify timing XML was injected
    timing = slide._element.find(qn("p:timing"))
    assert timing is not None


def test_apply_animations_unknown_effect():
    _, slide = _make_slide(1)
    animations = [{"target": "title", "effect": "nonexistent_effect"}]
    apply_animations(slide, animations)
    # Unknown effect should be skipped
    timing = slide._element.find(qn("p:timing"))
    assert timing is None


def test_apply_animations_empty_list():
    _, slide = _make_slide(1)
    apply_animations(slide, [])
    timing = slide._element.find(qn("p:timing"))
    assert timing is None


def test_apply_animations_no_matching_shapes():
    _, slide = _make_slide(1)
    animations = [{"target": "nonexistent_widget_xyz", "effect": "fade"}]
    apply_animations(slide, animations)
    timing = slide._element.find(qn("p:timing"))
    assert timing is None


def test_apply_animations_multiple_effects():
    _, slide = _make_slide(1)
    animations = [
        {"target": "title", "effect": "fade"},
        {"target": "content", "effect": "appear"},
    ]
    apply_animations(slide, animations)
    timing = slide._element.find(qn("p:timing"))
    assert timing is not None


def test_build_visual_effect_xml_unknown():
    """Unknown visual returns empty string."""
    xml = _build_visual_effect_xml("5", "unknown_effect_xyz")
    assert xml == ""


def test_build_click_par_contains_preset_id():
    effect = EFFECT_MAP["wipe"]
    xml = _build_click_par("3", effect)
    assert f'presetID="{effect["preset"]}"' in xml


# ---------------------------------------------------------------------------
# apply_animations
# ---------------------------------------------------------------------------


def test_apply_animations_adds_timing_element():
    _, slide = _make_slide(1)
    slide.shapes.title.text = "Test"
    anims = [{"target": "title", "effect": "fade"}]
    apply_animations(slide, anims)
    timing = slide._element.find(qn("p:timing"))
    assert timing is not None


def test_apply_animations_unique_ctn_ids():
    _, slide = _make_slide(1)
    slide.shapes.title.text = "Test"
    anims = [
        {"target": "title", "effect": "fade"},
        {"target": "content", "effect": "appear"},
    ]
    apply_animations(slide, anims)
    timing = slide._element.find(qn("p:timing"))
    ids = [el.get("id") for el in timing.iter(qn("p:cTn"))]
    assert len(ids) == len(set(ids)), "cTn IDs must be unique"


def test_apply_animations_empty_list_is_noop():
    _, slide = _make_slide(1)
    apply_animations(slide, [])
    timing = slide._element.find(qn("p:timing"))
    assert timing is None


def test_apply_animations_unknown_effect_warns(capsys):
    _, slide = _make_slide(1)
    anims = [{"target": "title", "effect": "teleport"}]
    apply_animations(slide, anims)
    captured = capsys.readouterr()
    assert "unknown animation" in captured.out.lower()


def test_apply_animations_replaces_existing_timing():
    _, slide = _make_slide(1)
    slide.shapes.title.text = "Test"
    anims = [{"target": "title", "effect": "fade"}]
    apply_animations(slide, anims)
    apply_animations(slide, anims)
    # Should have exactly one <p:timing>
    timings = slide._element.findall(qn("p:timing"))
    assert len(timings) == 1
