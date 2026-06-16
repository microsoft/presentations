"""Tests for src/slides.py."""

from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from presentations.slides import (
    SLIDE_BUILDERS,
    _CONTENT_DEFAULTS,
    _apply_position,
    _fit_badge_font_size,
    _hex_to_rgb,
    _interpolate_colors,
    _is_url,
    _scaled_defaults,
    _set_text_with_breaks,
    add_content_slide,
    add_resource_box_slide,
    add_section_header_slide,
    add_title_slide,
    add_two_column_slide,
)
from presentations.style import Style


@pytest.fixture()
def prs():
    return PptxPresentation()


@pytest.fixture()
def wide_prs():
    """A widescreen (16:9, 13.333" wide) presentation, as a rich layout forces."""
    p = PptxPresentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


@pytest.fixture()
def style():
    return Style()


# ---------------------------------------------------------------------------
# SLIDE_BUILDERS registry
# ---------------------------------------------------------------------------


class TestSlideBuilders:
    def test_known_types_registered(self):
        for t in ("title", "content", "section-header", "two-column", "resource-box"):
            assert t in SLIDE_BUILDERS

    def test_all_values_callable(self):
        for name, fn in SLIDE_BUILDERS.items():
            assert callable(fn), f"{name} is not callable"


# ---------------------------------------------------------------------------
# _apply_position
# ---------------------------------------------------------------------------


class TestApplyPosition:
    def test_sets_all_values(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shape = slide.shapes.title
        _apply_position(shape, {"left": 1.0, "top": 2.0, "width": 5.0, "height": 1.5})
        assert shape.left == Inches(1.0)
        assert shape.top == Inches(2.0)
        assert shape.width == Inches(5.0)
        assert shape.height == Inches(1.5)

    def test_partial_values(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shape = slide.shapes.title
        original_width = shape.width
        _apply_position(shape, {"left": 0.5, "top": 0.5})
        assert shape.left == Inches(0.5)
        assert shape.width == original_width  # unchanged

    def test_none_pos_is_noop(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shape = slide.shapes.title
        original = (shape.left, shape.top, shape.width, shape.height)
        _apply_position(shape, None)
        assert (shape.left, shape.top, shape.width, shape.height) == original


# ---------------------------------------------------------------------------
# _hex_to_rgb
# ---------------------------------------------------------------------------


class TestHexToRgb:
    def test_black(self):
        c = _hex_to_rgb("#000000")
        assert (c[0], c[1], c[2]) == (0, 0, 0)

    def test_white(self):
        c = _hex_to_rgb("#FFFFFF")
        assert (c[0], c[1], c[2]) == (255, 255, 255)

    def test_azure_blue(self):
        c = _hex_to_rgb("#0078D4")
        assert (c[0], c[1], c[2]) == (0, 120, 212)

    def test_no_hash(self):
        c = _hex_to_rgb("FF0000")
        assert (c[0], c[1], c[2]) == (255, 0, 0)


# ---------------------------------------------------------------------------
# _interpolate_colors
# ---------------------------------------------------------------------------


class TestInterpolateColors:
    def test_single_color(self):
        result = _interpolate_colors(["#FF0000"], 1)
        assert len(result) == 1

    def test_two_colors_two_slots(self):
        result = _interpolate_colors(["#000000", "#FFFFFF"], 2)
        assert len(result) == 2
        assert (result[0][0], result[0][1], result[0][2]) == (0, 0, 0)
        assert (result[1][0], result[1][1], result[1][2]) == (255, 255, 255)

    def test_zero_count(self):
        assert _interpolate_colors(["#FF0000"], 0) == []

    def test_gradient_midpoint(self):
        result = _interpolate_colors(["#000000", "#FFFFFF"], 3)
        # Midpoint should be ~(127, 127, 127)
        mid = result[1]
        assert 126 <= mid[0] <= 128


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


class TestAddTitleSlide:
    def test_creates_slide(self, prs, style):
        data = {
            "type": "title",
            "title": "Hello",
            "subtitle": "World",
            "notes": "Speaker notes",
            "animations": [],
            "positions": {},
        }
        add_title_slide(prs, data, style)
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Hello"
        assert slide.placeholders[1].text == "World"
        assert slide.notes_slide.notes_text_frame.text == "Speaker notes"

    def test_empty_subtitle(self, prs, style):
        data = {
            "type": "title",
            "title": "Only Title",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_title_slide(prs, data, style)
        assert prs.slides[0].placeholders[1].text == ""


class TestAddContentSlide:
    def test_creates_slide_with_bullets(self, prs, style):
        data = {
            "type": "content",
            "title": "Topics",
            "bullets": ["A", "B", "C"],
            "notes": "Notes here",
            "animations": [],
            "positions": {},
        }
        add_content_slide(prs, data, style)
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Topics"
        body = slide.placeholders[1].text_frame
        texts = [p.text for p in body.paragraphs]
        assert texts == ["A", "B", "C"]

    def test_empty_bullets(self, prs, style):
        data = {
            "type": "content",
            "title": "Empty",
            "bullets": [],
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_content_slide(prs, data, style)
        assert len(prs.slides) == 1


class TestAddSectionHeaderSlide:
    def test_creates_slide(self, prs, style):
        data = {
            "type": "section-header",
            "title": "Section",
            "subtitle": "Break time",
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_section_header_slide(prs, data, style)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Section"

    def test_no_subtitle(self, prs, style):
        data = {
            "type": "section-header",
            "title": "Part 2",
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_section_header_slide(prs, data, style)
        assert len(prs.slides) == 1


class TestAddTwoColumnSlide:
    def test_creates_slide(self, prs, style):
        data = {
            "type": "two-column",
            "title": "Compare",
            "left_bullets": ["L1", "L2"],
            "right_bullets": ["R1", "R2"],
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_two_column_slide(prs, data, style)
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        left_texts = [p.text for p in slide.placeholders[1].text_frame.paragraphs]
        right_texts = [p.text for p in slide.placeholders[2].text_frame.paragraphs]
        assert left_texts == ["L1", "L2"]
        assert right_texts == ["R1", "R2"]

    def test_empty_columns(self, prs, style):
        data = {
            "type": "two-column",
            "title": "Empty",
            "left_bullets": [],
            "right_bullets": [],
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_two_column_slide(prs, data, style)
        assert len(prs.slides) == 1


class TestAnimationsCallback:
    def test_animations_called(self, prs, style):
        called = []
        data = {
            "type": "title",
            "title": "Test",
            "subtitle": "",
            "notes": "",
            "animations": [{"target": "title", "effect": "fade"}],
            "positions": {},
        }
        add_title_slide(prs, data, style, apply_animations=lambda s, a: called.append((s, a)))
        assert len(called) == 1

    def test_no_animations_not_called(self, prs, style):
        called = []
        data = {
            "type": "title",
            "title": "Test",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_title_slide(prs, data, style, apply_animations=lambda s, a: called.append(1))
        assert len(called) == 0


# ---------------------------------------------------------------------------
# _is_url
# ---------------------------------------------------------------------------


class TestIsUrl:
    def test_http_url(self):
        assert _is_url("http://example.com") is True

    def test_https_url(self):
        assert _is_url("https://example.com/path?q=1") is True

    def test_plain_text(self):
        assert _is_url("just some text") is False

    def test_url_with_spaces_around(self):
        assert _is_url("  https://example.com  ") is True

    def test_text_with_space_not_url(self):
        assert _is_url("https://example.com is great") is False

    def test_empty_string(self):
        assert _is_url("") is False

    def test_ftp_not_matched(self):
        assert _is_url("ftp://files.example.com") is False


# ---------------------------------------------------------------------------
# _set_text_with_breaks
# ---------------------------------------------------------------------------


class TestSetTextWithBreaks:
    def test_single_line(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shape = slide.shapes.title
        _set_text_with_breaks(shape, "Hello World", Pt(20))
        assert shape.text_frame.paragraphs[0].text == "Hello World"

    def test_multiple_breaks(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shape = slide.shapes.title
        _set_text_with_breaks(shape, "Line1<br>Line2<br>Line3", Pt(18))
        texts = [p.text for p in shape.text_frame.paragraphs]
        assert texts == ["Line1", "Line2", "Line3"]

    def test_empty_text(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shape = slide.shapes.title
        _set_text_with_breaks(shape, "", Pt(20))
        assert shape.text_frame.paragraphs[0].text == ""


# ---------------------------------------------------------------------------
# add_title_slide – subtitle with <br>
# ---------------------------------------------------------------------------


class TestTitleSlideBreaks:
    def test_subtitle_with_br(self, prs, style):
        data = {
            "type": "title",
            "title": "Main Title",
            "subtitle": "Line A<br>Line B<br>Line C",
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_title_slide(prs, data, style)
        paras = prs.slides[0].placeholders[1].text_frame.paragraphs
        texts = [p.text for p in paras]
        assert texts == ["Line A", "Line B", "Line C"]


# ---------------------------------------------------------------------------
# add_content_slide – image width constraint
# ---------------------------------------------------------------------------


class TestContentSlideImageConstraint:
    def test_content_width_constrained_by_image(self, prs, style):
        """When an image is present and no explicit content position, width shrinks."""
        data = {
            "type": "content",
            "title": "With Image",
            "bullets": ["Bullet"],
            "image": {"path": "nonexistent.png", "left": 6.0, "top": 1.5},
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_content_slide(prs, data, style)
        body_ph = prs.slides[0].placeholders[1]
        expected_width = Inches(6.0 - 0.2 - 0.5)
        assert body_ph.width == expected_width


# ---------------------------------------------------------------------------
# add_resource_box_slide
# ---------------------------------------------------------------------------


class TestAddResourceBoxSlide:
    def test_creates_slide_with_boxes(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Resources",
            "subtitle": "Links",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [
                {
                    "label": "Learn",
                    "rows": [
                        {"name": "Docs", "url": "https://example.com"},
                    ],
                },
            ],
            "slide_style": {},
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1
        # Should have multiple shapes: textbox for title, subtitle, container, badge, etc.
        assert len(prs.slides[0].shapes) > 3

    def test_multiple_boxes(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Multi",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [
                {"label": "A", "rows": [{"name": "R1", "url": "http://a.com"}]},
                {"label": "B", "rows": [{"name": "R2", "url": "http://b.com"}]},
            ],
            "slide_style": {},
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1

    def test_empty_boxes(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Empty",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [],
            "slide_style": {},
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1

    def test_subtitle_with_gradient_colors(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Styled",
            "subtitle": "ABC",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [],
            "slide_style": {"SubtitleColors": "#FF0000, #00FF00, #0000FF"},
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1

    def test_slide_style_overrides(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Override",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [{"label": "X", "rows": [{"name": "N", "url": "http://x.com"}]}],
            "slide_style": {
                "SlideBackground": "#000000",
                "TitleColor": "#FFFFFF",
                "BoxBackground": "#333333",
            },
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1

    def test_outer_border(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Border",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [{"label": "B", "rows": [{"name": "N", "url": ""}]}],
            "slide_style": {"OuterBorderColor": "#FF0000"},
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1

    def test_multiple_rows_in_box(self, prs, style):
        data = {
            "type": "resource-box",
            "title": "Rows",
            "subtitle": "",
            "notes": "",
            "animations": [],
            "positions": {},
            "boxes": [
                {
                    "label": "Links",
                    "rows": [
                        {"name": "A", "url": "http://a.com"},
                        {"name": "B", "url": "http://b.com"},
                        {"name": "C", "url": "http://c.com"},
                    ],
                },
            ],
            "slide_style": {},
        }
        add_resource_box_slide(prs, data, style)
        assert len(prs.slides) == 1

    @staticmethod
    def _container_and_badge(slide, label):
        """Return (container_shape, badge_shape) for the box labelled *label*."""
        container = next(
            s for s in slide.shapes
            if str(s.shape_type) == "AUTO_SHAPE (1)"
            and not (s.has_text_frame and s.text_frame.text)
        )
        badge = next(
            s for s in slide.shapes
            if s.has_text_frame and s.text_frame.text == label
        )
        return container, badge

    def _build_single_box(self, prs, style, label, nrows=1):
        rows = [{"name": f"R{i}", "url": "https://example.com"} for i in range(nrows)]
        data = {
            "type": "resource-box", "title": "T", "subtitle": "", "notes": "",
            "animations": [], "positions": {},
            "boxes": [{"label": label, "rows": rows}], "slide_style": {},
        }
        add_resource_box_slide(prs, data, style)
        return self._container_and_badge(prs.slides[0], label)

    def test_badge_within_single_row_container(self, prs, style):
        # Issue #10: single-row badge must not protrude above/below its container.
        container, badge = self._build_single_box(prs, style, "Learn", nrows=1)
        ct, cb = container.top, container.top + container.height
        bt, bb = badge.top, badge.top + badge.height
        assert bt >= ct
        assert bb <= cb

    def test_badge_keeps_full_height_for_multi_row(self, prs, style):
        # Regression: taller multi-row containers should not shrink the badge.
        _, badge = self._build_single_box(prs, style, "Learn", nrows=3)
        assert badge.height == Inches(style.badge_height)

    def test_long_single_word_label_shrinks_font(self, prs, style):
        # Issue #11: an over-wide single word must scale down to avoid mid-word wrap.
        _, badge = self._build_single_box(prs, style, "Presentations", nrows=1)
        assert badge.text_frame.paragraphs[0].font.size < Pt(style.badge_font_size)

    def test_short_label_keeps_base_font(self, prs, style):
        _, badge = self._build_single_box(prs, style, "Learn", nrows=1)
        assert badge.text_frame.paragraphs[0].font.size == Pt(style.badge_font_size)

    def test_badge_autosize_enabled(self, prs, style):
        # Issue #11: PowerPoint shrink-on-overflow backstop is requested.
        _, badge = self._build_single_box(prs, style, "Presentations", nrows=1)
        assert badge.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


# ---------------------------------------------------------------------------
# _scaled_defaults – canvas-width scaling for core builders (issue #5)
# ---------------------------------------------------------------------------


class TestScaledDefaults:
    def test_unchanged_on_4_3_canvas(self, prs):
        """On a default 10" (4:3) deck the scale factor is exactly 1.0."""
        result = _scaled_defaults(prs, _CONTENT_DEFAULTS)
        assert result == _CONTENT_DEFAULTS
        # A fresh dict is returned (not the module-level constant)
        assert result is not _CONTENT_DEFAULTS

    def test_scales_horizontal_on_widescreen(self, wide_prs):
        """On a 13.333" deck, left/width scale by ~1.3333; top/height untouched."""
        scale = wide_prs.slide_width / Inches(10)
        result = _scaled_defaults(wide_prs, _CONTENT_DEFAULTS)
        title = result["title"]
        assert title["left"] == pytest.approx(0.5 * scale)
        assert title["width"] == pytest.approx(8.5 * scale)
        # Vertical values are constant across aspect ratios
        assert title["top"] == _CONTENT_DEFAULTS["title"]["top"]
        assert title["height"] == _CONTENT_DEFAULTS["title"]["height"]

    def test_does_not_mutate_input(self, wide_prs):
        before = {k: dict(v) for k, v in _CONTENT_DEFAULTS.items()}
        _scaled_defaults(wide_prs, _CONTENT_DEFAULTS)
        assert _CONTENT_DEFAULTS == before


class TestWidescreenCoreGeometry:
    """Regression for issue #5: core slides must span a forced 16:9 canvas."""

    def test_content_title_spans_widescreen(self, wide_prs, style):
        data = {
            "type": "content",
            "title": "Core Slide On The Right Canvas",
            "bullets": ["a"],
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_content_slide(wide_prs, data, style)
        title = wide_prs.slides[0].shapes.title
        scale = wide_prs.slide_width / Inches(10)
        # Title width scales to ~11.33" (8.5 x 1.3333), not the old 8.5"
        assert title.width == Inches(8.5 * scale)
        # Right-hand gap is the 4:3 gap (1.0") scaled, ~1.33" — not the ~4.33" bug
        right_gap = wide_prs.slide_width - title.left - title.width
        assert right_gap < Inches(1.5)

    def test_two_column_right_spans_widescreen(self, wide_prs, style):
        data = {
            "type": "two-column",
            "title": "Compare",
            "left_bullets": ["L"],
            "right_bullets": ["R"],
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_two_column_slide(wide_prs, data, style)
        scale = wide_prs.slide_width / Inches(10)
        right_ph = wide_prs.slides[0].placeholders[2]
        assert right_ph.left == Inches(5.0 * scale)
        assert right_ph.width == Inches(4.25 * scale)

    def test_section_header_subtitle_spans_widescreen(self, wide_prs, style):
        data = {
            "type": "section-header",
            "title": "Section",
            "subtitle": "Break",
            "notes": "",
            "animations": [],
            "positions": {},
        }
        add_section_header_slide(wide_prs, data, style)
        scale = wide_prs.slide_width / Inches(10)
        assert wide_prs.slides[0].shapes.title.width == Inches(9.0 * scale)


# ---------------------------------------------------------------------------
# _fit_badge_font_size
# ---------------------------------------------------------------------------


class TestFitBadgeFontSize:
    def test_short_word_keeps_base(self):
        assert _fit_badge_font_size("Learn", 0.8, 11) == 11

    def test_long_word_shrinks(self):
        assert _fit_badge_font_size("Presentations", 0.8, 11) < 11

    def test_never_below_floor(self):
        assert _fit_badge_font_size("Supercalifragilisticexpialidocious", 0.8, 11) >= 6

    def test_multi_word_uses_longest_word(self):
        # "Govern" (6) fits at base; the spaces should not force a shrink.
        assert _fit_badge_font_size("Build and Govern", 0.8, 11) == 11

    def test_empty_label_returns_base(self):
        assert _fit_badge_font_size("", 0.8, 11) == 11
