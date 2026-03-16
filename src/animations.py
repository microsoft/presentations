"""Animation engine – injects Open XML timing trees into slides.

All animation is built as click-triggered entrance effects using the standard
PowerPoint Open XML ``<p:timing>`` element.
"""

from __future__ import annotations

from lxml import etree
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Effect catalogue
# ---------------------------------------------------------------------------

# preset ID, preset class, subtype, visual-effect builder key
EFFECT_MAP: dict[str, dict[str, str]] = {
    "appear":       {"preset": "1",  "cls": "entr", "subtype": "0",  "visual": "none"},
    "fade":         {"preset": "10", "cls": "entr", "subtype": "0",  "visual": "fade"},
    "fly-in":       {"preset": "2",  "cls": "entr", "subtype": "4",  "visual": "fly-bottom"},
    "fly-in-left":  {"preset": "2",  "cls": "entr", "subtype": "8",  "visual": "fly-left"},
    "fly-in-right": {"preset": "2",  "cls": "entr", "subtype": "2",  "visual": "fly-right"},
    "fly-in-top":   {"preset": "2",  "cls": "entr", "subtype": "1",  "visual": "fly-top"},
    "wipe":         {"preset": "22", "cls": "entr", "subtype": "4",  "visual": "wipe"},
    "zoom":         {"preset": "23", "cls": "entr", "subtype": "0",  "visual": "zoom"},
    "float-in":     {"preset": "42", "cls": "entr", "subtype": "4",  "visual": "float"},
    "split":        {"preset": "16", "cls": "entr", "subtype": "0",  "visual": "split"},
    "blinds":       {"preset": "3",  "cls": "entr", "subtype": "0",  "visual": "blinds"},
}

# Target keyword → shape-name patterns
TARGET_SHAPE_MAP: dict[str, list[str]] = {
    "title":    ["Title"],
    "content":  ["Content Placeholder", "Text Placeholder"],
    "left":     ["Content Placeholder 2"],
    "right":    ["Content Placeholder 3"],
    "image":    ["Picture"],
    "subtitle": ["Subtitle"],
}

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---------------------------------------------------------------------------
# Shape resolution
# ---------------------------------------------------------------------------


def _resolve_target_shapes(slide, target: str) -> list:
    """Resolve an English target keyword to actual shape objects on the slide."""
    if target == "all":
        return list(slide.shapes)
    if target == "bullets":
        target = "content"
    patterns = TARGET_SHAPE_MAP.get(target, [])
    matched = []
    for shape in slide.shapes:
        for pat in patterns:
            if pat.lower() in shape.name.lower():
                matched.append(shape)
                break
    if not matched:
        for shape in slide.shapes:
            if target.lower() in shape.name.lower():
                matched.append(shape)
    return matched


def _get_shape_id(shape) -> str:
    """Get the numeric ``id`` attribute from a shape's XML element."""
    sp_elem = shape._element
    for nv_tag in ("p:nvSpPr", "p:nvPicPr", "p:nvGrpSpPr", "p:nvCxnSpPr"):
        nv = sp_elem.find(qn(nv_tag))
        if nv is not None:
            cNvPr = nv.find(qn("p:cNvPr"))
            if cNvPr is not None:
                return cNvPr.get("id")
    return "2"


# ---------------------------------------------------------------------------
# Visual-effect XML builders
# ---------------------------------------------------------------------------


def _build_visual_effect_xml(shape_id: str, visual: str) -> str:  # noqa: C901
    """Return extra animation XML nodes for a specific visual effect."""
    if visual == "none":
        return ""
    if visual == "fade":
        return (
            f'<p:animEffect transition="in" filter="fade">'
            f"<p:cBhvr>"
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"</p:cBhvr>"
            f"</p:animEffect>"
        )
    if visual.startswith("fly-"):
        direction = visual.split("-", 1)[1]
        if direction == "bottom":
            attr, fr, to = "ppt_y", "#ppt_h+#ppt_y", "#ppt_y"
        elif direction == "top":
            attr, fr, to = "ppt_y", "-#ppt_h", "#ppt_y"
        elif direction == "left":
            attr, fr, to = "ppt_x", "-#ppt_w", "#ppt_x"
        elif direction == "right":
            attr, fr, to = "ppt_x", "#ppt_w+#ppt_x", "#ppt_x"
        else:
            attr, fr, to = "ppt_y", "#ppt_h+#ppt_y", "#ppt_y"
        return (
            f'<p:anim calcmode="lin" valueType="num">'
            f'<p:cBhvr additive="base">'
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"<p:attrNameLst><p:attrName>{attr}</p:attrName></p:attrNameLst>"
            f"</p:cBhvr>"
            f"<p:tavLst>"
            f'<p:tav tm="0"><p:val><p:strVal val="{fr}"/></p:val></p:tav>'
            f'<p:tav tm="100000"><p:val><p:strVal val="{to}"/></p:val></p:tav>'
            f"</p:tavLst>"
            f"</p:anim>"
        )
    if visual == "wipe":
        return (
            f'<p:animEffect transition="in" filter="wipe(down)">'
            f"<p:cBhvr>"
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"</p:cBhvr>"
            f"</p:animEffect>"
        )
    if visual == "zoom":
        return (
            f"<p:animScale>"
            f"<p:cBhvr>"
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"</p:cBhvr>"
            f'<p:by x="0" y="0"/>'
            f'<p:from x="0" y="0"/>'
            f'<p:to x="100000" y="100000"/>'
            f"</p:animScale>"
        )
    if visual == "float":
        return (
            f'<p:anim calcmode="lin" valueType="num">'
            f'<p:cBhvr additive="base">'
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"<p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst>"
            f"</p:cBhvr>"
            f"<p:tavLst>"
            f'<p:tav tm="0"><p:val><p:strVal val="#ppt_y+0.1"/></p:val></p:tav>'
            f'<p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>'
            f"</p:tavLst>"
            f"</p:anim>"
            f'<p:animEffect transition="in" filter="fade">'
            f"<p:cBhvr>"
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"</p:cBhvr>"
            f"</p:animEffect>"
        )
    if visual == "split":
        return (
            f'<p:animEffect transition="in" filter="barn(inVertical)">'
            f"<p:cBhvr>"
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"</p:cBhvr>"
            f"</p:animEffect>"
        )
    if visual == "blinds":
        return (
            f'<p:animEffect transition="in" filter="blinds(horizontal)">'
            f"<p:cBhvr>"
            f'<p:cTn id="0" dur="500" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f"</p:cBhvr>"
            f"</p:animEffect>"
        )
    return ""


# ---------------------------------------------------------------------------
# Click-par builder
# ---------------------------------------------------------------------------


def _build_click_par(shape_id: str, effect: dict) -> str:
    """Build one click-triggered animation group (3-level ``<p:par>`` nesting)."""
    preset = effect["preset"]
    cls = effect["cls"]
    subtype = effect["subtype"]
    visual_xml = _build_visual_effect_xml(shape_id, effect["visual"])

    return (
        f'<p:par xmlns:p="{_P_NS}" xmlns:a="{_A_NS}">'
        f'<p:cTn id="0" fill="hold">'
        f"<p:stCondLst>"
        f'<p:cond delay="indefinite"/>'
        f"</p:stCondLst>"
        f"<p:childTnLst>"
        f"<p:par>"
        f'<p:cTn id="0" fill="hold">'
        f"<p:stCondLst>"
        f'<p:cond delay="0"/>'
        f"</p:stCondLst>"
        f"<p:childTnLst>"
        f"<p:par>"
        f'<p:cTn id="0" presetID="{preset}" presetClass="{cls}" presetSubtype="{subtype}"'
        f' fill="hold" nodeType="clickEffect">'
        f"<p:stCondLst>"
        f'<p:cond delay="0"/>'
        f"</p:stCondLst>"
        f"<p:childTnLst>"
        f"<p:set>"
        f"<p:cBhvr>"
        f'<p:cTn id="0" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f"</p:cTn>"
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f"<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>"
        f"</p:cBhvr>"
        f'<p:to><p:strVal val="visible"/></p:to>'
        f"</p:set>"
        f"{visual_xml}"
        f"</p:childTnLst>"
        f"</p:cTn>"
        f"</p:par>"
        f"</p:childTnLst>"
        f"</p:cTn>"
        f"</p:par>"
        f"</p:childTnLst>"
        f"</p:cTn>"
        f"</p:par>"
    )


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def apply_animations(slide, animations: list[dict]) -> None:
    """Inject animation XML into *slide* for each animation spec entry."""
    if not animations:
        return

    anim_pairs: list[tuple[str, dict]] = []
    for anim in animations:
        effect_name = anim["effect"]
        effect = EFFECT_MAP.get(effect_name)
        if effect is None:
            print(
                f"Warning: unknown animation '{effect_name}', skipping. "
                f"Available: {', '.join(EFFECT_MAP.keys())}"
            )
            continue
        shapes = _resolve_target_shapes(slide, anim["target"])
        if not shapes:
            print(
                f"Warning: no shapes matched target '{anim['target']}' on slide, "
                f"skipping animation."
            )
            continue
        for s in shapes:
            anim_pairs.append((_get_shape_id(s), effect))

    if not anim_pairs:
        return

    click_pars = "".join(_build_click_par(sid, eff) for sid, eff in anim_pairs)

    timing_xml = (
        f'<p:timing xmlns:p="{_P_NS}" xmlns:a="{_A_NS}">'
        f"<p:tnLst>"
        f"<p:par>"
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f"<p:childTnLst>"
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f"<p:childTnLst>"
        f"{click_pars}"
        f"</p:childTnLst>"
        f"</p:cTn>"
        f"<p:prevCondLst>"
        f'<p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
        f"</p:prevCondLst>"
        f"<p:nextCondLst>"
        f'<p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
        f"</p:nextCondLst>"
        f"</p:seq>"
        f"</p:childTnLst>"
        f"</p:cTn>"
        f"</p:par>"
        f"</p:tnLst>"
        f"</p:timing>"
    )

    timing_el = etree.fromstring(timing_xml.encode())

    # Assign unique sequential IDs to all cTn elements
    ctn_id = 1
    for ctn in timing_el.iter(qn("p:cTn")):
        ctn.set("id", str(ctn_id))
        ctn_id += 1

    # Remove any existing <p:timing> and append ours
    existing = slide._element.find(qn("p:timing"))
    if existing is not None:
        slide._element.remove(existing)
    slide._element.append(timing_el)
